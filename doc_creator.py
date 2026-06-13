"""
Hujjat yaratuvchi modul - To'liq Python (node.js yo'q)
python-pptx va python-docx kutubxonalari
"""

import os
import asyncio
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DocPt, RGBColor as DocRGB, Inches as DocInches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import tempfile

logger = logging.getLogger(__name__)

OUTPUT_DIR = "/tmp/edubot_files"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Rang palitrasisi
PRIMARY   = RGBColor(0x1A, 0x3C, 0x5E)
SECONDARY = RGBColor(0x2D, 0x7D, 0xD2)
ACCENT    = RGBColor(0xF0, 0xA5, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xEA, 0xF2, 0xFB)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY      = RGBColor(0x99, 0x99, 0x99)


def _add_colored_background(slide, prs, color: RGBColor):
    """Slaydi rang bilan to'ldirish"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, prs, x, y, w, h, color: RGBColor):
    """To'rtburchak shakl qo'shish"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    return shape


def _add_textbox(slide, text, x, y, w, h,
                  font_size=18, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


async def create_pptx(topic: str, content: dict) -> str:
    """python-pptx bilan prezentatsiya yaratish"""
    safe_name = "".join(c for c in topic[:30] if c.isalnum() or c in " _-").strip()
    output_path = os.path.join(OUTPUT_DIR, f"pptx_{safe_name}_{datetime.now().strftime('%H%M%S')}.pptx")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout
    title_str = content.get("title", topic)
    sections  = content.get("sections", [])

    # ── 1. SARLAVHA SLAYDI ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    _add_colored_background(slide, prs, PRIMARY)
    _add_rect(slide, prs, 0, 0, 13.33, 0.12, ACCENT)
    _add_rect(slide, prs, 0, 7.38, 13.33, 0.12, ACCENT)

    _add_textbox(slide, title_str, 1, 1.8, 11.33, 2.5,
                  font_size=36, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
    _add_textbox(slide, f"Mavzu: {topic[:80]}", 1, 4.5, 11.33, 0.6,
                  font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)
    today = datetime.now().strftime("%Y-yil, %B")
    _add_textbox(slide, today, 1, 6.5, 11.33, 0.5,
                  font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # ── 2. MUNDARIJA ────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank_layout)
    _add_colored_background(slide2, prs, WHITE)
    _add_rect(slide2, prs, 0, 0, 13.33, 1.2, PRIMARY)
    _add_textbox(slide2, "📋  Mundarija", 0.5, 0.2, 12, 0.8,
                  font_size=28, bold=True, color=WHITE)

    toc_text = "\n".join(f"  {i+1}.  {s['heading']}" for i, s in enumerate(sections))
    _add_textbox(slide2, toc_text, 0.8, 1.4, 11.5, 5.5,
                  font_size=16, color=DARK_TEXT)

    # ── 3. KONTENT SLAYDLARI ────────────────────────────────────────
    icons = ["💡","📌","🔬","📊","📖","⭐","🎯","🔑"]
    for idx, section in enumerate(sections):
        heading = section.get("heading", f"Bo'lim {idx+1}")
        raw    = section.get("content", "")
        parts  = [p.strip() for p in raw.split("||") if p.strip()] or [raw[:300]]

        bg = LIGHT if idx % 2 == 0 else WHITE
        sl = prs.slides.add_slide(blank_layout)
        _add_colored_background(sl, prs, bg)
        _add_rect(sl, prs, 0, 0, 13.33, 1.2, PRIMARY)

        icon = icons[idx % len(icons)]
        _add_textbox(sl, f"{icon}  {heading}", 0.5, 0.18, 12, 0.85,
                      font_size=24, bold=True, color=WHITE)

        bullet_text = "\n\n".join(f"  •  {p}" for p in parts[:5])
        _add_textbox(sl, bullet_text, 0.6, 1.4, 12.1, 5.7,
                      font_size=15, color=DARK_TEXT)

        _add_rect(sl, prs, 0, 7.3, 13.33, 0.1, SECONDARY)
        num_box = _add_textbox(sl, str(idx + 2), 12.6, 7.1, 0.5, 0.35,
                                font_size=10, color=GRAY, align=PP_ALIGN.CENTER)

    # ── 4. YAKUNIY SLAYD ────────────────────────────────────────────
    end = prs.slides.add_slide(blank_layout)
    _add_colored_background(end, prs, PRIMARY)
    _add_rect(end, prs, 0, 0, 13.33, 0.12, ACCENT)
    _add_rect(end, prs, 0, 7.38, 13.33, 0.12, ACCENT)
    _add_textbox(end, "E'tiboringiz uchun\nrahmat! 🙏", 1.5, 2.2, 10, 3,
                  font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    logger.info(f"✅ PPTX tayyor: {output_path}")
    return output_path


async def create_docx(topic: str, content: dict, doc_type: str) -> str:
    """python-docx bilan Word hujjat yaratish"""
    safe_name = "".join(c for c in topic[:30] if c.isalnum() or c in " _-").strip()
    output_path = os.path.join(OUTPUT_DIR, f"docx_{safe_name}_{datetime.now().strftime('%H%M%S')}.docx")

    doc = Document()

    # ── Sahifa o'lchamlari (A4) ──────────────────────────────────────
    section = doc.sections[0]
    section.page_width  = Cm(21)
    section.page_height = Cm(29.7)
    section.left_margin   = Cm(3)
    section.right_margin  = Cm(1.5)
    section.top_margin    = Cm(2)
    section.bottom_margin = Cm(2)

    doc_type_names = {
        "mustaqil": "MUSTAQIL ISH",
        "amaliy":   "AMALIY ISH",
        "referat":  "REFERAT",
    }
    doc_type_name = doc_type_names.get(doc_type, "HUJJAT")
    title_str = content.get("title", topic)
    sections  = content.get("sections", [])

    # ── Sarlavha sahifasi ────────────────────────────────────────────
    def center_bold(text, size=14, color=None, space_before=0, space_after=6):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        pf = p.paragraph_format
        pf.space_before = DocPt(space_before)
        pf.space_after  = DocPt(space_after)
        run = p.add_run(text)
        run.bold = True
        run.font.size = DocPt(size)
        run.font.name = "Times New Roman"
        if color:
            run.font.color.rgb = DocRGB(*color)
        return p

    center_bold("O'ZBEKISTON RESPUBLIKASI OLIY TA'LIM VAZIRLIGI", 12, space_before=40)
    center_bold(doc_type_name, 18, color=(0x1A, 0x3C, 0x5E), space_before=30, space_after=10)

    p_mavzu = doc.add_paragraph()
    p_mavzu.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r1 = p_mavzu.add_run("Mavzu: ")
    r1.font.size = DocPt(14)
    r1.font.name = "Times New Roman"
    r2 = p_mavzu.add_run(title_str)
    r2.bold = True
    r2.font.size = DocPt(14)
    r2.font.name = "Times New Roman"

    center_bold(f"{datetime.now().year}-yil", 12, space_before=80)

    doc.add_page_break()

    # ── Asosiy kontent ───────────────────────────────────────────────
    for section_data in sections:
        heading = section_data.get("heading", "Bo'lim")
        body    = section_data.get("content", "")

        # Sarlavha
        h = doc.add_heading(heading, level=1)
        h.alignment = WD_ALIGN_PARAGRAPH.LEFT
        for run in h.runs:
            run.font.name  = "Times New Roman"
            run.font.size  = DocPt(14)
            run.font.color.rgb = DocRGB(0x1A, 0x3C, 0x5E)

        # Body — paragraflar
        for para_text in body.split("\n"):
            para_text = para_text.strip()
            if not para_text:
                continue
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            pf = p.paragraph_format
            pf.first_line_indent = Cm(1.25)
            pf.space_after  = DocPt(6)
            pf.line_spacing = DocPt(18)
            run = p.add_run(para_text)
            run.font.size = DocPt(14)
            run.font.name = "Times New Roman"

    doc.save(output_path)
    logger.info(f"✅ DOCX tayyor: {output_path}")
    return output_path
